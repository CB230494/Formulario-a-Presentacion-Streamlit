import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
import datetime

st.set_page_config(page_title="Informe Policial Automático", layout="centered")

st.title("🚔 Generador de Informe de Dispositivo Policial")

# ---- FORMULARIO ----
with st.form("formulario_informe"):
    st.subheader("🔹 Datos del Informe")
    nombre_dispositivo = st.text_input("Nombre del Dispositivo")
    nombre_responsable = st.text_input("Nombre del Responsable")
    cargo_responsable = st.text_input("Cargo del Responsable")
    direccion_regional = st.selectbox("Dirección Regional", [
        "Cartago", "San José", "Alajuela", "Heredia", "Guanacaste", "Puntarenas", "Limón"
    ])
    delegacion_policial = st.text_input("Delegación Policial")
    fecha_ejecucion = st.date_input("Fecha de ejecución")

    st.subheader("🔹 Contenido del Informe")
    descripcion_resultados = st.text_area("Breve descripción de resultados obtenidos")
    analisis_operativo = st.text_area("Análisis o balance operativo")
    recomendaciones = st.text_area("Recomendaciones o sugerencias")

    enviar = st.form_submit_button("📤 Generar Informe PPTX")

# ---- FUNCIONES ----
@st.cache_resource
def cargar_plantilla(url):
    respuesta = requests.get(url)
    return BytesIO(respuesta.content)

def reemplazar_portada(prs, delegacion, direccion):
    portada = prs.slides[0]
    for shape in portada.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "DELEGACION_POLICIAL" in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = delegacion
                p.font.bold = True
                p.font.size = Pt(48)
                p.font.name = 'Arial'
            if "DIRECCION_REGIONAL" in text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = direccion
                p.font.bold = True
                p.font.size = Pt(36)
                p.font.name = 'Arial'

def crear_diapositiva_con_estilo(prs, titulo_texto, contenido_texto):
    slide_layout = prs.slide_layouts[6]  # Slide vacío
    slide = prs.slides.add_slide(slide_layout)

    # Fondo azul claro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)

    # Título
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = titulo_texto
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 51, 102)  # Azul oscuro

    # Contenido
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    p = tf.add_paragraph()
    p.text = contenido_texto
    p.font.size = Pt(24)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(0, 0, 0)
    tf.word_wrap = True

def generar_pptx(datos, plantilla_bytes):
    prs = Presentation(plantilla_bytes)

    # Guardar portada (slide 0) y página institucional (última slide)
    portada = prs.slides[0]
    pagina_institucional = prs.slides[1]

    # Eliminar todo excepto la portada
    for i in range(len(prs.slides) - 1, 0, -1):
        r_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[i]

    # Reemplazar textos en la portada
    reemplazar_portada(prs, datos['delegacion_policial'], datos['direccion_regional'])

    # Insertar nuevas diapositivas generadas
    crear_diapositiva_con_estilo(prs, "Información General", 
        f"Nombre del Dispositivo: {datos['nombre_dispositivo']}\n"
        f"Responsable: {datos['nombre_responsable']}\n"
        f"Cargo del Responsable: {datos['cargo_responsable']}\n"
        f"Fecha de Ejecución: {datos['fecha_ejecucion']}"
    )

    crear_diapositiva_con_estilo(prs, "Resultados Obtenidos", datos['descripcion_resultados'])
    crear_diapositiva_con_estilo(prs, "Análisis Operativo", datos['analisis_operativo'])
    crear_diapositiva_con_estilo(prs, "Recomendaciones", datos['recomendaciones'])

    # Insertar la página institucional al final
    slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in pagina_institucional.shapes:
        new_element = shape.element
        new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# ---- DESPUÉS DE ENVIAR FORMULARIO ----
if enviar:
    campos = [
        nombre_dispositivo, nombre_responsable, cargo_responsable,
        direccion_regional, delegacion_policial, fecha_ejecucion,
        descripcion_resultados, analisis_operativo, recomendaciones
    ]
    if not all(campos):
        st.error("⚠️ Completa todos los campos para generar el informe.")
    else:
        st.success("✅ Informe generado correctamente.")

        plantilla_url = "https://github.com/CB230494/Formulario-a-Presentacion-Streamlit/raw/refs/heads/main/plantilla_personalizada.pptx"
        plantilla_bytes = cargar_plantilla(plantilla_url)

        datos = {
            'nombre_dispositivo': nombre_dispositivo,
            'nombre_responsable': nombre_responsable,
            'cargo_responsable': cargo_responsable,
            'direccion_regional': f"Dirección Regional {direccion_regional}",
            'delegacion_policial': delegacion_policial,
            'fecha_ejecucion': fecha_ejecucion.strftime("%d/%m/%Y"),
            'descripcion_resultados': descripcion_resultados,
            'analisis_operativo': analisis_operativo,
            'recomendaciones': recomendaciones
        }

        ppt_buffer = generar_pptx(datos, plantilla_bytes)

        nombre_archivo = f"Informe_{nombre_dispositivo.replace(' ', '_')}_{datetime.datetime.now().strftime('%Y%m%d')}.pptx"

        st.download_button(
            label="📥 Descargar Informe PPTX",
            data=ppt_buffer,
            file_name=nombre_archivo,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

